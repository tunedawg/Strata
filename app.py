"""
Universal Boolean Search Tool
------------------------------
Supported file types:
  PDF (text + OCR fallback), DOCX, DOC, PPTX, CSV, XLSX, XLS, TXT, MD,
  EML (email), MSG (Outlook email)

Install dependencies:
  pip install flask pdfplumber pypdf python-docx openpyxl reportlab
  pip install pytesseract pdf2image pillow python-pptx extract-msg

Also install Tesseract OCR engine:
  Windows: https://github.com/UB-Mannheim/tesseract/wiki
  Mac:     brew install tesseract
  Linux:   sudo apt install tesseract-ocr
"""

import os, re, sys, json, pickle, threading, time, difflib, email
from datetime import timedelta
from pathlib import Path
from collections import defaultdict

from flask import Flask, render_template, request, jsonify, send_file

app = Flask(__name__)

# ── CONFIG ────────────────────────────────────────────────────────────────────
# When packaged with PyInstaller, UNIVERSAL_SEARCH_DATA is set by launcher.py
# so data lives in ~/Documents/UniversalSearch rather than next to the binary.
BASE_DIR     = os.environ.get(
    "UNIVERSAL_SEARCH_DATA",
    os.path.join(os.path.expanduser("~"), "universal_search_data"),
)
DATASETS_DIR = os.path.join(BASE_DIR, "datasets")
EXPORTS_DIR  = os.path.join(BASE_DIR, "exports")
os.makedirs(DATASETS_DIR, exist_ok=True)
os.makedirs(EXPORTS_DIR,  exist_ok=True)

# OCR settings
OCR_ENABLED   = True   # set False to skip OCR entirely
OCR_DPI       = 200    # higher = better quality but slower
OCR_MIN_CHARS = 100    # if a page has fewer chars than this, run OCR on it

# ── Portable Tesseract path resolution ───────────────────────────────────────
# Works whether running as plain Python or a frozen PyInstaller bundle.
# Priority: (1) bundled copy next to the exe, (2) standard install locations,
# (3) let pytesseract find it on PATH automatically.
def _find_tesseract():
    import shutil

    # 1. Bundled alongside the frozen exe (developers can ship tesseract/tesseract.exe)
    if hasattr(sys, "_MEIPASS"):
        bundled = os.path.join(sys._MEIPASS, "tesseract", "tesseract.exe")
        if os.path.isfile(bundled):
            return bundled

    # 2. Standard Windows install locations
    if sys.platform == "win32":
        candidates = [
            r"C:\Program Files\Tesseract-OCR\tesseract.exe",
            r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
            os.path.join(os.environ.get("LOCALAPPDATA", ""), r"Programs\Tesseract-OCR\tesseract.exe"),
            os.path.join(os.environ.get("APPDATA", ""), r"Programs\Tesseract-OCR\tesseract.exe"),
        ]
        for c in candidates:
            if os.path.isfile(c):
                return c

    # 3. On PATH (macOS / Linux / Windows if added to PATH)
    found = shutil.which("tesseract")
    if found:
        return found

    return None   # OCR will be disabled gracefully

_TESSERACT_CMD = _find_tesseract()
try:
    import pytesseract
    if _TESSERACT_CMD:
        pytesseract.pytesseract.tesseract_cmd = _TESSERACT_CMD
except ImportError:
    pass

# ── Portable Poppler path resolution ─────────────────────────────────────────
# pdf2image needs Poppler binaries on Windows (Mac/Linux usually have it via brew/apt).
def _find_poppler():
    import shutil

    # 1. Bundled alongside the frozen exe (developers can ship poppler/bin/)
    if hasattr(sys, "_MEIPASS"):
        bundled = os.path.join(sys._MEIPASS, "poppler", "bin")
        if os.path.isdir(bundled):
            return bundled

    # 2. On PATH already (macOS via brew, Linux via apt, or Windows if added)
    if shutil.which("pdftoppm"):
        return None   # already on PATH — pdf2image finds it automatically

    # 3. Common Windows install locations
    if sys.platform == "win32":
        candidates = [
            r"C:\Program Files\poppler\bin",
            r"C:\Program Files (x86)\poppler\bin",
            os.path.join(os.environ.get("LOCALAPPDATA", ""), r"Programs\poppler\bin"),
            os.path.join(os.path.expanduser("~"), r"Documents\poppler\bin"),
        ]
        for c in candidates:
            if os.path.isdir(c):
                return c

    return None   # pdf2image will try its own detection

POPPLER_PATH = _find_poppler()

# In-memory state
LOADED_DATASET  = None
INDEX           = {}
INDEX_READY     = False
INDEX_PROGRESS  = {"status": "idle", "done": 0, "total": 0, "dataset": ""}
INDEX_LOCK      = threading.Lock()


def ocr_available():
    try:
        import pytesseract
        if _TESSERACT_CMD:
            pytesseract.pytesseract.tesseract_cmd = _TESSERACT_CMD
        pytesseract.get_tesseract_version()
        return True
    except Exception:
        return False

_ocr_ready = None  # cached check

def run_ocr_on_image(img):
    """Run OCR on a PIL Image and return text string."""
    try:
        import pytesseract
        if _TESSERACT_CMD:
            pytesseract.pytesseract.tesseract_cmd = _TESSERACT_CMD
        return pytesseract.image_to_string(img, lang="eng")
    except Exception:
        return ""

def pdf_page_to_image(path, page_num, dpi=OCR_DPI):
    """Convert a single PDF page to a PIL Image for OCR."""
    try:
        from pdf2image import convert_from_path
        images = convert_from_path(path, dpi=dpi, first_page=page_num+1,
                                   last_page=page_num+1,
                                   poppler_path=POPPLER_PATH)
        return images[0] if images else None
    except Exception:
        return None


# ── CHUNKING / EXTRACTION ─────────────────────────────────────────────────────

def extract_pdf(path):
    """Extract text from PDF. Falls back to OCR for pages with little/no text."""
    global _ocr_ready
    import pdfplumber

    chunks = {}
    if _ocr_ready is None:
        _ocr_ready = ocr_available() and OCR_ENABLED
        print(f"  OCR ready: {_ocr_ready}")

    try:
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = (page.extract_text() or "").strip()
                print(f"  Page {i+1}: {len(text)} chars from pdfplumber")

                # If page has very little text, try OCR
                if len(text) < OCR_MIN_CHARS and _ocr_ready:
                    print(f"  Page {i+1}: trying OCR...")
                    img = pdf_page_to_image(path, i)
                    if img:
                        ocr_text = run_ocr_on_image(img).strip()
                        print(f"  Page {i+1}: OCR got {len(ocr_text)} chars")
                        if len(ocr_text) > len(text):
                            text = ocr_text
                    else:
                        print(f"  Page {i+1}: pdf_page_to_image returned None")

                if text.strip():
                    chunks[f"p{i+1}"] = text.lower()
    except Exception as e:
        print(f"  PDF error {path}: {e}")

    if not chunks and _ocr_ready:
        print(f"  No chunks found, trying full OCR...")
        try:
            from pdf2image import convert_from_path
            images = convert_from_path(path, dpi=OCR_DPI, poppler_path=POPPLER_PATH)
            for i, img in enumerate(images):
                text = run_ocr_on_image(img).strip()
                if text:
                    chunks[f"p{i+1}"] = text.lower()
        except Exception as e:
            print(f"  PDF full-OCR error {path}: {e}")

    return chunks

def extract_docx(path):
    """Extract text from Word documents including tables."""
    from docx import Document
    chunks = {}
    try:
        doc  = Document(path)
        buf  = []
        idx  = 1

        # Paragraphs
        for para in doc.paragraphs:
            t = para.text.strip()
            if t:
                buf.append(t)
            if len(buf) >= 5:
                chunks[f"para{idx}"] = " ".join(buf).lower()
                idx += 1
                buf = []

        # Tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    buf.append(row_text)
                if len(buf) >= 5:
                    chunks[f"para{idx}"] = " ".join(buf).lower()
                    idx += 1
                    buf = []

        if buf:
            chunks[f"para{idx}"] = " ".join(buf).lower()
    except Exception as e:
        print(f"  DOCX error {path}: {e}")
    return chunks

def extract_pptx(path):
    """Extract text from PowerPoint slides."""
    chunks = {}
    try:
        from pptx import Presentation
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
                # Tables in slides
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                texts.append(cell.text.strip())
            if texts:
                chunks[f"slide{i+1}"] = " ".join(texts).lower()
    except Exception as e:
        print(f"  PPTX error {path}: {e}")
    return chunks

def extract_csv(path):
    """Extract text from CSV files."""
    import csv
    chunks = {}
    try:
        with open(path, encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            for i, row in enumerate(reader):
                text = " ".join(str(c) for c in row).strip()
                if text:
                    chunks[f"row{i+1}"] = text.lower()
    except Exception as e:
        print(f"  CSV error {path}: {e}")
    return chunks

def extract_xlsx(path):
    """Extract text from Excel files including all cell types."""
    import openpyxl
    chunks = {}
    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        try:
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    parts = []
                    for c in row:
                        if c is None: continue
                        # Handle dates, numbers, strings uniformly
                        try:
                            parts.append(str(c).strip())
                        except Exception:
                            pass
                    text = " ".join(p for p in parts if p)
                    if text:
                        chunks[f"{sheet}_row{i+1}"] = text.lower()
        finally:
            wb.close()
    except Exception as e:
        print(f"  XLSX error {path}: {e}")
    return chunks

def extract_txt(path):
    """Extract text from plain text and markdown files."""
    chunks = {}
    try:
        text  = Path(path).read_text(encoding="utf-8", errors="replace")
        paras = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
        for i, p in enumerate(paras):
            chunks[f"para{i+1}"] = p.lower()
    except Exception as e:
        print(f"  TXT error {path}: {e}")
    return chunks

def extract_eml(path):
    """Extract text from .eml email files."""
    chunks = {}
    try:
        with open(path, "rb") as f:
            msg = email.message_from_bytes(f.read())

        parts = []
        # Headers
        for header in ("From", "To", "Cc", "Subject", "Date"):
            val = msg.get(header, "")
            if val:
                parts.append(f"{header}: {val}")

        # Body
        if msg.is_multipart():
            for part in msg.walk():
                ct = part.get_content_type()
                if ct == "text/plain":
                    try:
                        body = part.get_payload(decode=True).decode("utf-8", errors="replace")
                        parts.append(body)
                    except Exception:
                        pass
        else:
            try:
                body = msg.get_payload(decode=True).decode("utf-8", errors="replace")
                parts.append(body)
            except Exception:
                pass

        full_text = "\n".join(parts)
        paras = [p.strip() for p in re.split(r"\n\s*\n", full_text) if p.strip()]
        for i, p in enumerate(paras):
            chunks[f"para{i+1}"] = p.lower()

    except Exception as e:
        print(f"  EML error {path}: {e}")
    return chunks

def extract_msg(path):
    """Extract text from Outlook .msg email files."""
    chunks = {}
    try:
        import extract_msg as em
        msg = em.Message(path)
        parts = []
        if msg.sender:    parts.append(f"From: {msg.sender}")
        if msg.to:        parts.append(f"To: {msg.to}")
        if msg.cc:        parts.append(f"Cc: {msg.cc}")
        if msg.subject:   parts.append(f"Subject: {msg.subject}")
        if msg.date:      parts.append(f"Date: {msg.date}")
        if msg.body:      parts.append(msg.body)

        full_text = "\n".join(parts)
        paras = [p.strip() for p in re.split(r"\n\s*\n", full_text) if p.strip()]
        for i, p in enumerate(paras):
            chunks[f"para{i+1}"] = p.lower()
        msg.close()
    except Exception as e:
        print(f"  MSG error {path}: {e}")
    return chunks

def extract_image(path):
    """Extract text from image files using OCR."""
    global _ocr_ready
    chunks = {}
    if _ocr_ready is None:
        _ocr_ready = ocr_available() and OCR_ENABLED
    if not _ocr_ready:
        return chunks
    try:
        from PIL import Image
        img  = Image.open(path)
        text = run_ocr_on_image(img).strip()
        if text:
            chunks["p1"] = text.lower()
    except Exception as e:
        print(f"  Image OCR error {path}: {e}")
    return chunks

def extract_tiff(path):
    """Extract text from multi-page TIFF files using OCR."""
    global _ocr_ready
    chunks = {}
    if _ocr_ready is None:
        _ocr_ready = ocr_available() and OCR_ENABLED
    if not _ocr_ready:
        return chunks
    try:
        from PIL import Image
        img = Image.open(path)
        page = 0
        while True:
            text = run_ocr_on_image(img).strip()
            if text:
                chunks[f"p{page+1}"] = text.lower()
            page += 1
            try:
                img.seek(page)
            except EOFError:
                break
    except Exception as e:
        print(f"  TIFF OCR error {path}: {e}")
    return chunks

EXTRACTORS = {
    # Documents
    ".pdf":  extract_pdf,
    ".docx": extract_docx,
    ".doc":  extract_docx,
    ".pptx": extract_pptx,
    ".ppt":  extract_pptx,
    # Spreadsheets
    ".csv":  extract_csv,
    ".xlsx": extract_xlsx,
    ".xls":  extract_xlsx,
    # Text
    ".txt":  extract_txt,
    ".md":   extract_txt,
    ".rtf":  extract_txt,
    # Email
    ".eml":  extract_eml,
    ".msg":  extract_msg,
    # Images (OCR)
    ".png":  extract_image,
    ".jpg":  extract_image,
    ".jpeg": extract_image,
    ".tif":  extract_tiff,
    ".tiff": extract_tiff,
    ".bmp":  extract_image,
}


# ── INDEXING ──────────────────────────────────────────────────────────────────

def index_path(dataset_name):
    return os.path.join(DATASETS_DIR, dataset_name, "_index.pkl")

def get_datasets():
    datasets = []
    for d in sorted(Path(DATASETS_DIR).iterdir()):
        if d.is_dir():
            idx  = index_path(d.name)
            info = {
                "name":    d.name,
                "indexed": os.path.exists(idx),
                "files":   len([f for f in d.rglob("*")
                                if f.is_file()
                                and not f.name.startswith("_")
                                and f.suffix.lower() in EXTRACTORS]),
            }
            datasets.append(info)
    return datasets

def build_index_for(dataset_name):
    global INDEX, INDEX_READY, INDEX_PROGRESS, LOADED_DATASET
    folder = Path(DATASETS_DIR) / dataset_name
    files  = [f for f in folder.rglob("*")
              if f.is_file() and not f.name.startswith("_")
              and f.suffix.lower() in EXTRACTORS]

    print(f"Starting index for '{dataset_name}': {len(files)} files found")

    with INDEX_LOCK:
        INDEX_PROGRESS = {"status": "indexing", "done": 0,
                          "total": len(files), "dataset": dataset_name}
        INDEX_READY    = False

    local_index = {}
    for i, fpath in enumerate(files):
        try:
            ext       = fpath.suffix.lower()
            extractor = EXTRACTORS[ext]
            chunks    = extractor(str(fpath))
            rel       = str(fpath.relative_to(folder))
            local_index[rel] = chunks
            with INDEX_LOCK:
                INDEX_PROGRESS["done"] = i + 1
            if (i + 1) % 50 == 0:
                print(f"  Indexed {i+1}/{len(files)}...")
        except Exception as e:
            print(f"  ERROR on file {fpath.name}: {e}")

    print(f"Saving index for '{dataset_name}'...")
    with open(index_path(dataset_name), "wb") as f:
        pickle.dump(local_index, f)

    with INDEX_LOCK:
        INDEX          = local_index
        INDEX_READY    = True
        LOADED_DATASET = dataset_name
        INDEX_PROGRESS["status"] = "ready"
    print(f"Done! Indexed {len(local_index)} files for '{dataset_name}'.")

def load_index_for(dataset_name):
    global INDEX, INDEX_READY, INDEX_PROGRESS, LOADED_DATASET
    with INDEX_LOCK:
        INDEX_PROGRESS = {"status": "loading", "done": 0, "total": 0, "dataset": dataset_name}
        INDEX_READY    = False
    with open(index_path(dataset_name), "rb") as f:
        local_index = pickle.load(f)
    with INDEX_LOCK:
        INDEX          = local_index
        INDEX_READY    = True
        LOADED_DATASET = dataset_name
        INDEX_PROGRESS = {"status": "ready", "done": len(local_index),
                          "total": len(local_index), "dataset": dataset_name}
    print(f"Loaded index for '{dataset_name}': {len(local_index)} files.")


# ── QUERY ENGINE ──────────────────────────────────────────────────────────────

def words_of(text):
    return re.findall(r"[a-z0-9']+", text.lower())

def match_term(text, token):
    token = token.strip()
    if token.startswith('"') and token.endswith('"'):
        return token[1:-1].lower() in text
    fuzzy = token.startswith("~") or token.endswith("~")
    token = token.strip("~")
    if "*" in token:
        pattern = re.compile(r"\b" + re.escape(token).replace(r"\*", r"\w*") + r"\b")
        return bool(pattern.search(text))
    if fuzzy:
        matches = difflib.get_close_matches(token.lower(), words_of(text), n=1, cutoff=0.75)
        return bool(matches)
    return token.lower() in text

def proximity_match(text, phrase_a, phrase_b, within):
    def find_positions(wlist, phrase):
        pw = words_of(phrase)
        n  = len(pw)
        return [i for i in range(len(wlist) - n + 1) if wlist[i:i+n] == pw]
    wlist = words_of(text)
    pos_a = find_positions(wlist, phrase_a.strip('"'))
    pos_b = find_positions(wlist, phrase_b.strip('"'))
    if not pos_a or not pos_b:
        return False
    return any(abs(a - b) <= within for a in pos_a for b in pos_b)

def parse_query(query):
    prox_re = re.compile(r'(".*?"|[\w~*]+)\s+(?:W|NEAR)/(\d+)\s+(".*?"|[\w~*]+)', re.IGNORECASE)

    def evaluate(text, q):
        q = q.strip()
        if not q:
            return False
        parts = split_top(q, " OR ")
        if len(parts) > 1:
            return any(evaluate(text, p) for p in parts)
        parts = split_top(q, " AND NOT ")
        if len(parts) > 1:
            return evaluate(text, parts[0]) and not any(evaluate(text, p) for p in parts[1:])
        parts = split_top(q, " AND ")
        if len(parts) > 1:
            return all(evaluate(text, p) for p in parts)
        # Infix NOT: "foo NOT bar" — foo must match, bar must not
        parts = split_top(q, " NOT ")
        if len(parts) > 1:
            return evaluate(text, parts[0]) and not any(evaluate(text, p) for p in parts[1:])
        if q.upper().startswith("NOT "):
            return not evaluate(text, q[4:])
        if q.startswith("(") and q.endswith(")") and matching_close(q) == len(q) - 1:
            return evaluate(text, q[1:-1])
        pm = prox_re.search(q)
        if pm:
            return proximity_match(text, pm.group(1), pm.group(3), int(pm.group(2)))
        return match_term(text, q)

    def split_top(q, op):
        parts, depth, in_quote, buf = [], 0, False, ""
        i = 0
        while i < len(q):
            c = q[i]
            if c == '"':
                in_quote = not in_quote; buf += c
            elif not in_quote and c == "(":
                depth += 1; buf += c
            elif not in_quote and c == ")":
                depth -= 1; buf += c
            elif not in_quote and depth == 0 and q[i:i+len(op)].upper() == op.upper():
                parts.append(buf.strip()); buf = ""; i += len(op); continue
            else:
                buf += c
            i += 1
        parts.append(buf.strip())
        return [p for p in parts if p]

    def matching_close(s):
        depth = 0
        for i, c in enumerate(s):
            if c == "(": depth += 1
            elif c == ")":
                depth -= 1
                if depth == 0: return i
        return -1

    return lambda text: evaluate(text, query)

def run_search_with_index(terms, index):
    results = {}
    for term in terms:
        term = term.strip()
        if not term:
            continue
        try:
            matcher = parse_query(term)
        except Exception as e:
            results[term] = {"error": str(e), "total_hits": 0, "docs": []}
            continue

        docs_hit = []
        total    = 0
        for doc_name, chunks in index.items():
            matched = [cid for cid, text in chunks.items() if matcher(text)]
            if matched:
                docs_hit.append({"name": doc_name, "chunks": matched})
                total += len(matched)

        docs_hit.sort(key=lambda x: -len(x["chunks"]))
        results[term] = {"total_hits": total, "docs": docs_hit}
    return results

def run_search(terms):
    return run_search_with_index(terms, INDEX)

def search_dataset(dataset_name, terms):
    """Load dataset pickle locally (no global state mutation) and search it."""
    pkl = index_path(dataset_name)
    if not os.path.exists(pkl):
        raise FileNotFoundError(f"Dataset '{dataset_name}' is not indexed")
    with open(pkl, "rb") as f:
        local_index = pickle.load(f)
    return run_search_with_index(terms, local_index)


# ── PDF EXPORT ────────────────────────────────────────────────────────────────

def export_results(term, docs_hit, dataset_name):
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors
    from reportlab.lib.units import inch

    safe     = re.sub(r'[\\/*?:"<>|]', "_", term)[:60]
    out_path = os.path.join(EXPORTS_DIR, f"search_{safe}.pdf")

    doc    = SimpleDocTemplate(out_path, pagesize=letter,
                               leftMargin=0.75*inch, rightMargin=0.75*inch,
                               topMargin=0.75*inch,  bottomMargin=0.75*inch)
    styles = getSampleStyleSheet()
    title_s = ParagraphStyle("T", parent=styles["Title"], fontSize=16, spaceAfter=4)
    sub_s   = ParagraphStyle("S", parent=styles["Normal"], fontSize=9,
                              textColor=colors.grey, spaceAfter=12)
    doc_s   = ParagraphStyle("D", parent=styles["Heading2"], fontSize=11,
                              textColor=colors.HexColor("#1a1a2e"), spaceAfter=4)
    chunk_s = ParagraphStyle("C", parent=styles["Normal"], fontSize=9,
                              leading=13, spaceAfter=8,
                              backColor=colors.HexColor("#f8f8f8"),
                              borderPadding=(6, 8, 6, 8))
    label_s = ParagraphStyle("L", parent=styles["Normal"], fontSize=7,
                              textColor=colors.grey, spaceAfter=2)

    def esc(t):
        return str(t).replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")

    story = [
        Paragraph(esc(f'Search: {term}'), title_s),
        Paragraph(esc(f'Dataset: {dataset_name}  ·  {sum(len(d["chunks"]) for d in docs_hit)} matching chunks across {len(docs_hit)} file(s)'), sub_s),
        HRFlowable(width="100%", thickness=1, color=colors.HexColor("#CCCCCC"), spaceAfter=12),
    ]

    for entry in docs_hit:
        doc_name = entry["name"]
        chunks   = entry["chunks"]
        story.append(Paragraph(esc(doc_name), doc_s))

        # Look up actual chunk texts from index
        doc_chunks = INDEX.get(doc_name, {})
        for cid in chunks:
            text = doc_chunks.get(cid, "")
            story.append(Paragraph(esc(cid), label_s))
            # Truncate very long chunks for readability
            display = text[:800] + ("…" if len(text) > 800 else "")
            story.append(Paragraph(esc(display), chunk_s))

        story.append(Spacer(1, 10))

    doc.build(story)
    return out_path


# ── HIT REPORT ────────────────────────────────────────────────────────────────

def generate_hit_report(results, dataset_name):
    """Generate both Excel and PDF hit reports. Returns (xlsx_path, pdf_path)"""
    from datetime import datetime as dt
    timestamp = dt.now().strftime("%Y%m%d_%H%M%S")
    now_str   = dt.now().strftime("%B %d, %Y %I:%M %p")
    xlsx_path = os.path.join(EXPORTS_DIR, f"hit_report_{timestamp}.xlsx")
    pdf_path  = os.path.join(EXPORTS_DIR, f"hit_report_{timestamp}.pdf")

    # ── Build data ────────────────────────────────────────────────────────────
    summary_rows  = []
    detailed_rows = []

    for term, r in results.items():
        if r.get("error"):
            summary_rows.append({
                "Search Term":    term,
                "Total Pages":    0,
                "Total Contacts": 0,
                "Contacts":       [f"ERROR: {r['error']}"],
            })
            continue
        docs  = r.get("docs", [])
        total = r.get("total_hits", 0)
        # Store as list to avoid splitting on commas in contact names
        contact_names = [d["name"].replace(".pdf", "") for d in docs]
        summary_rows.append({
            "Search Term":    term,
            "Total Pages":    total,
            "Total Contacts": len(docs),
            "Contacts":       contact_names,
        })
        for doc in docs:
            contact = doc["name"].replace(".pdf", "")
            for chunk in doc["chunks"]:
                page_match = re.match(r'^p(\d+)$', chunk)
                page_num   = int(page_match.group(1)) if page_match else chunk
                detailed_rows.append({
                    "Search Term": term,
                    "Contact":     contact,
                    "Page":        page_num,
                    "Chunk ID":    chunk,
                })

    # ── Excel ─────────────────────────────────────────────────────────────────
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    wb  = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "Summary"

    hdr_fill   = PatternFill("solid", fgColor="1A1A2E")
    hdr_font   = Font(bold=True, color="FFFFFF")
    # openpyxl requires 8-char ARGB hex for border colors
    hdr_border = Border(bottom=Side(style="medium", color="FF4F8EF7"))
    stripe     = PatternFill("solid", fgColor="F0F4FF")

    summary_headers = ["Search Term", "Total Pages", "Total Contacts", "Contacts"]
    for col, h in enumerate(summary_headers, 1):
        cell = ws1.cell(row=1, column=col, value=h)
        cell.font      = hdr_font
        cell.fill      = hdr_fill
        cell.border    = hdr_border
        cell.alignment = Alignment(horizontal="center")

    for row_idx, row in enumerate(summary_rows, 2):
        contacts_str = ", ".join(row["Contacts"][:10])
        if len(row["Contacts"]) > 10:
            contacts_str += f" (+{len(row['Contacts'])-10} more)"
        ws1.cell(row=row_idx, column=1, value=row["Search Term"])
        ws1.cell(row=row_idx, column=2, value=row["Total Pages"]).alignment  = Alignment(horizontal="center")
        ws1.cell(row=row_idx, column=3, value=row["Total Contacts"]).alignment = Alignment(horizontal="center")
        ws1.cell(row=row_idx, column=4, value=contacts_str)
        if row_idx % 2 == 0:
            for col in range(1, 5):
                ws1.cell(row=row_idx, column=col).fill = stripe

    ws1.column_dimensions["A"].width = 35
    ws1.column_dimensions["B"].width = 14
    ws1.column_dimensions["C"].width = 16
    ws1.column_dimensions["D"].width = 60

    ws2 = wb.create_sheet("Detailed")
    detail_headers = ["Search Term", "Contact", "Page", "Chunk ID"]
    for col, h in enumerate(detail_headers, 1):
        cell = ws2.cell(row=1, column=col, value=h)
        cell.font      = hdr_font
        cell.fill      = hdr_fill
        cell.border    = hdr_border
        cell.alignment = Alignment(horizontal="center")

    for row_idx, row in enumerate(detailed_rows, 2):
        ws2.cell(row=row_idx, column=1, value=row["Search Term"])
        ws2.cell(row=row_idx, column=2, value=row["Contact"])
        ws2.cell(row=row_idx, column=3, value=row["Page"]).alignment  = Alignment(horizontal="center")
        ws2.cell(row=row_idx, column=4, value=row["Chunk ID"]).alignment = Alignment(horizontal="center")
        if row_idx % 2 == 0:
            for col in range(1, 5):
                ws2.cell(row=row_idx, column=col).fill = stripe

    ws2.column_dimensions["A"].width = 35
    ws2.column_dimensions["B"].width = 35
    ws2.column_dimensions["C"].width = 10
    ws2.column_dimensions["D"].width = 14

    wb.save(xlsx_path)

    # ── PDF ───────────────────────────────────────────────────────────────────
    from reportlab.lib.pagesizes import letter, landscape
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                    HRFlowable, Table, TableStyle)
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors as rc
    from reportlab.lib.units import inch

    # Use unique style name prefix to avoid re-registration errors
    _pfx = f"HR{timestamp}_"
    base = getSampleStyleSheet()

    doc = SimpleDocTemplate(pdf_path, pagesize=landscape(letter),
                            leftMargin=0.75*inch, rightMargin=0.75*inch,
                            topMargin=0.75*inch,  bottomMargin=0.75*inch)

    DARK   = rc.HexColor("#1a1a2e")
    ACCENT = rc.HexColor("#4f8ef7")
    STRIPE = rc.HexColor("#f0f4ff")

    title_s = ParagraphStyle(_pfx+"title", parent=base["Title"],   fontSize=18, spaceAfter=4)
    sub_s   = ParagraphStyle(_pfx+"sub",   parent=base["Normal"],  fontSize=9,
                              textColor=rc.grey, spaceAfter=16)
    head_s  = ParagraphStyle(_pfx+"head",  parent=base["Heading1"],fontSize=13, spaceAfter=8)
    cell_s  = ParagraphStyle(_pfx+"cell",  parent=base["Normal"],  fontSize=8,  leading=11)

    def esc(t):
        return str(t).replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")

    total_hits = sum(r.get("total_hits", 0) for r in results.values())
    story = [
        Paragraph(esc(f"Hit Report — {dataset_name}"), title_s),
        Paragraph(esc(f"Generated: {now_str}  ·  {len(results)} term(s)  ·  "
                      f"{total_hits:,} total page hits"), sub_s),
        HRFlowable(width="100%", thickness=2, color=ACCENT, spaceAfter=16),
        Paragraph("Summary", head_s),
    ]

    # Summary table — guard against empty
    sum_data = [["Search Term", "Total Pages", "Total Contacts", "Contacts (first 5)"]]
    for row in summary_rows:
        contacts_short = ", ".join(row["Contacts"][:5])
        if len(row["Contacts"]) > 5:
            contacts_short += f" (+{len(row['Contacts'])-5} more)"
        sum_data.append([
            Paragraph(esc(row["Search Term"]), cell_s),
            str(row["Total Pages"]),
            str(row["Total Contacts"]),
            Paragraph(esc(contacts_short), cell_s),
        ])

    if len(sum_data) == 1:
        sum_data.append(["No results", "", "", ""])

    sum_table = Table(sum_data, colWidths=[2.8*inch, 1*inch, 1.2*inch, 4.5*inch])
    sum_table.setStyle(TableStyle([
        ("BACKGROUND",     (0,0), (-1,0), DARK),
        ("TEXTCOLOR",      (0,0), (-1,0), rc.white),
        ("FONTNAME",       (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTSIZE",       (0,0), (-1,0), 9),
        ("ALIGN",          (1,0), (2,-1), "CENTER"),
        ("FONTSIZE",       (0,1), (-1,-1), 8),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [rc.white, STRIPE]),
        ("GRID",           (0,0), (-1,-1), 0.25, rc.HexColor("#cccccc")),
        ("LINEBELOW",      (0,0), (-1,0), 1.5, ACCENT),
        ("VALIGN",         (0,0), (-1,-1), "MIDDLE"),
        ("TOPPADDING",     (0,0), (-1,-1), 4),
        ("BOTTOMPADDING",  (0,0), (-1,-1), 4),
    ]))
    story += [sum_table, Spacer(1, 20), Paragraph("Detailed Hits", head_s)]

    # Detailed table — guard against empty
    det_data = [["Search Term", "Contact", "Page", "Chunk"]]
    for row in detailed_rows:
        det_data.append([
            Paragraph(esc(row["Search Term"]), cell_s),
            Paragraph(esc(row["Contact"]), cell_s),
            str(row["Page"]),
            str(row["Chunk ID"]),
        ])

    if len(det_data) == 1:
        det_data.append(["No detailed hits", "", "", ""])

    det_table = Table(det_data, colWidths=[2.8*inch, 3.5*inch, 0.8*inch, 0.8*inch])
    det_table.setStyle(TableStyle([
        ("BACKGROUND",     (0,0), (-1,0), DARK),
        ("TEXTCOLOR",      (0,0), (-1,0), rc.white),
        ("FONTNAME",       (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTSIZE",       (0,0), (-1,0), 9),
        ("ALIGN",          (2,0), (3,-1), "CENTER"),
        ("FONTSIZE",       (0,1), (-1,-1), 8),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [rc.white, STRIPE]),
        ("GRID",           (0,0), (-1,-1), 0.25, rc.HexColor("#cccccc")),
        ("LINEBELOW",      (0,0), (-1,0), 1.5, ACCENT),
        ("VALIGN",         (0,0), (-1,-1), "MIDDLE"),
        ("TOPPADDING",     (0,0), (-1,-1), 4),
        ("BOTTOMPADDING",  (0,0), (-1,-1), 4),
    ]))
    story.append(det_table)

    doc.build(story)
    return xlsx_path, pdf_path

@app.route("/hit_report", methods=["POST"])
def hit_report():
    if not INDEX_READY:
        return jsonify({"error": "No dataset loaded"}), 503
    data    = request.json
    results = data.get("results", {})
    fmt     = data.get("format", "xlsx")
    dataset = LOADED_DATASET
    try:
        xlsx_path, pdf_path = generate_hit_report(results, dataset)
        if fmt == "pdf":
            return send_file(pdf_path, as_attachment=True,
                             download_name=os.path.basename(pdf_path))
        else:
            return send_file(xlsx_path, as_attachment=True,
                             download_name=os.path.basename(xlsx_path))
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.after_request
def no_cache(response):
    response.headers["Cache-Control"] = "no-store, no-cache, must-revalidate, max-age=0"
    response.headers["Pragma"]        = "no-cache"
    response.headers["Expires"]       = "0"
    return response

@app.route("/")
def index():
    return render_template("index.html")

@app.route("/datasets", methods=["GET"])
def list_datasets():
    return jsonify(get_datasets())

@app.route("/datasets/create", methods=["POST"])
def create_dataset():
    name = request.json.get("name", "").strip()
    name = re.sub(r'[\\/*?:"<>|]', "_", name)
    if not name:
        return jsonify({"error": "Name required"}), 400
    path = Path(DATASETS_DIR) / name
    if path.exists():
        return jsonify({"error": "Dataset already exists"}), 400
    path.mkdir(parents=True)
    return jsonify({"ok": True, "name": name})

@app.route("/datasets/<n>/delete", methods=["POST"])
def delete_dataset(n):
    import shutil
    folder = Path(DATASETS_DIR) / n
    if not folder.exists():
        return jsonify({"error": "Dataset not found"}), 404
    try:
        shutil.rmtree(str(folder))
        global INDEX, INDEX_READY, LOADED_DATASET, INDEX_PROGRESS
        if LOADED_DATASET == n:
            with INDEX_LOCK:
                INDEX          = {}
                INDEX_READY    = False
                LOADED_DATASET = None
                INDEX_PROGRESS = {"status": "idle", "done": 0, "total": 0, "dataset": ""}
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/datasets/<name>/upload", methods=["POST"])
def upload_files(name):
    folder = Path(DATASETS_DIR) / name
    if not folder.exists():
        return jsonify({"error": "Dataset not found"}), 404
    saved = []
    for f in request.files.getlist("files"):
        # Normalise separators and strip leading slashes
        raw = (f.filename or "").replace("\\", "/").strip("/")
        # Remove any ".." path traversal components
        parts = [p for p in raw.split("/") if p and p != ".."]
        if not parts:
            continue
        dest = folder.joinpath(*parts)
        # Hard safety check: dest must stay inside the dataset folder
        try:
            dest.resolve().relative_to(folder.resolve())
        except ValueError:
            continue
        dest.parent.mkdir(parents=True, exist_ok=True)
        f.save(str(dest))
        saved.append("/".join(parts))
    return jsonify({"ok": True, "saved": saved})

@app.route("/datasets/<name>/index", methods=["POST"])
def trigger_index(name):
    folder = Path(DATASETS_DIR) / name
    if not folder.exists():
        return jsonify({"error": "Dataset not found"}), 404
    t = threading.Thread(target=build_index_for, args=(name,), daemon=True)
    t.start()
    return jsonify({"ok": True})

@app.route("/datasets/<name>/load", methods=["POST"])
def load_dataset(name):
    idx = index_path(name)
    if not os.path.exists(idx):
        return jsonify({"error": "Not indexed yet"}), 400
    t = threading.Thread(target=load_index_for, args=(name,), daemon=True)
    t.start()
    return jsonify({"ok": True})

@app.route("/index_status")
def index_status():
    with INDEX_LOCK:
        return jsonify({**INDEX_PROGRESS, "ready": INDEX_READY,
                        "loaded": LOADED_DATASET})

@app.route("/search", methods=["POST"])
def search():
    if not INDEX_READY:
        return jsonify({"error": "No dataset loaded"}), 503
    terms = request.json.get("terms", [])
    return jsonify({"results": run_search(terms), "dataset": LOADED_DATASET})

@app.route("/export", methods=["POST"])
def export():
    if not INDEX_READY:
        return jsonify({"error": "No dataset loaded"}), 503
    data    = request.json
    term    = data.get("term", "")
    docs    = data.get("docs", [])
    dataset = LOADED_DATASET
    try:
        path = export_results(term, docs, dataset)
        return send_file(path, as_attachment=True, download_name=os.path.basename(path))
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/preview/<dataset>/<path:doc_name>")
def preview_file(dataset, doc_name):
    """Serve a PDF file directly so the browser can open it."""
    folder   = Path(DATASETS_DIR) / dataset
    filepath = folder / doc_name
    # Security check — make sure path stays inside the dataset folder
    try:
        filepath.resolve().relative_to(folder.resolve())
    except ValueError:
        return "Forbidden", 403
    if not filepath.exists():
        return "File not found", 404
    # For PDFs, serve inline so browser opens the viewer
    ext = filepath.suffix.lower()
    if ext == ".pdf":
        return send_file(str(filepath), mimetype="application/pdf")
    else:
        return send_file(str(filepath), as_attachment=True)

@app.route("/production")
def production():
    return render_template("production.html")

@app.route("/production/search", methods=["POST"])
def production_search():
    data    = request.json or {}
    dataset = data.get("dataset", "").strip()
    terms   = data.get("terms", [])
    if not dataset:
        return jsonify({"error": "dataset required"}), 400
    try:
        results = search_dataset(dataset, terms)
    except FileNotFoundError as e:
        return jsonify({"error": str(e)}), 404
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    return jsonify({"dataset": dataset, "results": results})

@app.route("/production/export_zip", methods=["POST"])
def production_export_zip():
    import zipfile
    import io
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment

    data       = request.json or {}
    dataset    = data.get("dataset", "").strip()
    selections = data.get("selections", [])

    if not dataset:
        return jsonify({"error": "dataset required"}), 400
    if not selections:
        return jsonify({"error": "no documents selected"}), 400

    folder = Path(DATASETS_DIR) / dataset
    if not folder.exists():
        return jsonify({"error": "Dataset not found"}), 404

    # ── Responsiveness log ────────────────────────────────────────────────────
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Responsiveness Log"

    hdr_font = Font(bold=True, color="FFFFFF")
    hdr_fill = PatternFill("solid", fgColor="1A1A2E")
    stripe   = PatternFill("solid", fgColor="F0F4FF")
    headers  = ["Filename", "Responsive Terms", "Pages / Chunks"]
    for col, h in enumerate(headers, 1):
        cell           = ws.cell(row=1, column=col, value=h)
        cell.font      = hdr_font
        cell.fill      = hdr_fill
        cell.alignment = Alignment(horizontal="center")
    ws.column_dimensions["A"].width = 50
    ws.column_dimensions["B"].width = 55
    ws.column_dimensions["C"].width = 35

    for row_idx, sel in enumerate(selections, 2):
        filename   = sel.get("name", "").replace("\\", "/")
        terms_str  = ", ".join(sel.get("terms", []))
        chunks_str = ", ".join(sel.get("chunks", []))
        ws.cell(row=row_idx, column=1, value=filename)
        ws.cell(row=row_idx, column=2, value=terms_str)
        ws.cell(row=row_idx, column=3, value=chunks_str)
        if row_idx % 2 == 0:
            for col in range(1, 4):
                ws.cell(row=row_idx, column=col).fill = stripe

    xlsx_buf = io.BytesIO()
    wb.save(xlsx_buf)
    xlsx_buf.seek(0)

    # ── ZIP ───────────────────────────────────────────────────────────────────
    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, mode="w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("responsiveness_log.xlsx", xlsx_buf.read())
        for sel in selections:
            rel_path = sel.get("name", "")
            src_path = folder / rel_path
            try:
                src_path.resolve().relative_to(folder.resolve())
            except ValueError:
                continue
            if src_path.exists():
                arc_name = "files/" + rel_path.replace("\\", "/")
                zf.write(str(src_path), arcname=arc_name)

    zip_buf.seek(0)
    return send_file(zip_buf, mimetype="application/zip", as_attachment=True,
                     download_name="production_export.zip")


# ── REDACTION SCANNER STATE ───────────────────────────────────────────────────
REDACT_CANCEL   = False
REDACT_PROGRESS = {"status": "idle", "done": 0, "total": 0, "current": ""}
REDACT_RESULTS  = []
REDACT_LOCK     = threading.Lock()


# ── REDACTION DETECTION HELPERS ───────────────────────────────────────────────

def _color_is_black(c):
    """Return True if color value represents black/near-black."""
    if c is None:
        return False
    try:
        if isinstance(c, (int, float)):
            return float(c) < 0.15
        if isinstance(c, (list, tuple)):
            if len(c) == 1:
                return float(c[0]) < 0.15
            if len(c) == 3:    # RGB 0–1
                return all(float(v) < 0.15 for v in c)
            if len(c) == 4:    # CMYK: high K = black
                return float(c[3]) > 0.85
    except Exception:
        pass
    return False


def _color_is_white(c):
    """Return True if color value represents white/near-white.
    Tightened to avoid false positives on light-colored speech bubbles
    (e.g. light green RGB ~0.93,0.99,0.92) in screenshot-based PDFs.
    """
    if c is None:
        return False
    try:
        if isinstance(c, (int, float)):
            return float(c) > 0.97
        if isinstance(c, (list, tuple)):
            if len(c) == 1:
                return float(c[0]) > 0.97
            if len(c) == 3:    # RGB 0–1 — all channels must be near 1.0
                vals = [float(v) for v in c]
                # Must be near-white: all high AND low variance (no tints)
                return all(v > 0.95 for v in vals) and (max(vals) - min(vals)) < 0.04
            if len(c) == 4:    # CMYK: (0,0,0,0) = white
                return all(float(v) < 0.04 for v in c)
    except Exception:
        pass
    return False


def detect_redactions_in_pdf(path):
    """
    Multi-method PDF redaction detection.
    Returns list of finding dicts: {page, type, location, details}
    """
    findings = []

    # Method 1: official PDF /Redact annotation objects
    try:
        import pypdf
        reader = pypdf.PdfReader(path, strict=False)
        for page_num, page in enumerate(reader.pages, 1):
            annots = page.get('/Annots')
            if not annots:
                continue
            for ref in annots:
                try:
                    a = ref.get_object() if hasattr(ref, 'get_object') else ref
                    if str(a.get('/Subtype', '')) == '/Redact':
                        rect = a.get('/Rect', [])
                        try:
                            coords = [round(float(x), 1) for x in rect]
                        except Exception:
                            coords = list(rect)
                        overlay = str(a.get('/OverlayText', '')).strip()
                        findings.append({
                            'page': page_num,
                            'type': 'PDF Redaction Annotation',
                            'location': f"Page {page_num} — coords {coords}",
                            'details': (f"Standard /Redact annotation"
                                        + (f"; overlay: '{overlay}'" if overlay else '')),
                        })
                except Exception:
                    pass
    except Exception:
        pass

    # Method 2: pdfplumber — black bars, white boxes, white-ink text
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                pw = float(page.width  or 612)
                ph = float(page.height or 792)

                # --- rectangles ---
                for rect in (page.rects or []):
                    try:
                        x0  = float(rect.get('x0',  0))
                        x1  = float(rect.get('x1',  0))
                        top = float(rect.get('top', 0))
                        bot = float(rect.get('bottom', top))
                        w   = abs(x1 - x0)
                        h   = abs(bot - top)
                        if w < 15 or h < 3:
                            continue
                        fill_color = rect.get('non_stroking_color')
                        if fill_color is None:
                            continue

                        if _color_is_black(fill_color) and h <= 72:
                            # Narrow black rectangle — classic redaction bar
                            # Exclude thin lines (h < 4) and full-page borders
                            if h < 4 or (w / pw > 0.95 and h / ph > 0.95):
                                continue
                            findings.append({
                                'page': page_num,
                                'type': 'Black Redaction Bar',
                                'location': f"Page {page_num} — x:{x0:.0f}–{x1:.0f}, y:{top:.0f}–{bot:.0f}",
                                'details': f"{w:.0f}×{h:.0f}pt black-filled rectangle",
                            })
                        # White box / whiteout detection removed — too many false positives
                    except Exception:
                        pass

                # White text detection removed — too many false positives
    except Exception:
        pass

    return findings


def detect_redactions_in_docx(path):
    """
    Detect white/hidden text in Word .docx files.
    Returns list of finding dicts: {page, type, location, details}
    """
    findings = []
    try:
        from docx import Document
        from docx.oxml.ns import qn
        doc = Document(path)

        for para_idx, para in enumerate(doc.paragraphs, 1):
            for run in para.runs:
                preview = (run.text or '').strip()[:50].replace('\n', ' ')
                if not preview:
                    continue

                # White font color
                try:
                    clr = run.font.color
                    if clr and clr.type is not None and clr.rgb is not None:
                        r = int(str(clr.rgb)[0:2], 16)
                        g = int(str(clr.rgb)[2:4], 16)
                        b = int(str(clr.rgb)[4:6], 16)
                        if r > 240 and g > 240 and b > 240:
                            findings.append({
                                'page': None,
                                'type': 'White Font Color (DOCX)',
                                'location': f"Paragraph {para_idx}",
                                'details': f"Text in white: '{preview}'",
                            })
                except Exception:
                    pass

                # Hidden text
                try:
                    if run.font.hidden:
                        findings.append({
                            'page': None,
                            'type': 'Hidden Text (DOCX)',
                            'location': f"Paragraph {para_idx}",
                            'details': f"Hidden text run: '{preview}'",
                        })
                except Exception:
                    pass

                # White highlight
                try:
                    rpr = run._r.find(qn('w:rPr'))
                    if rpr is not None:
                        hl = rpr.find(qn('w:highlight'))
                        if hl is not None and hl.get(qn('w:val'), '') == 'white':
                            findings.append({
                                'page': None,
                                'type': 'White Highlight (DOCX)',
                                'location': f"Paragraph {para_idx}",
                                'details': f"White-highlighted text: '{preview}'",
                            })
                except Exception:
                    pass

    except Exception as e:
        print(f"  DOCX redaction scan error {path}: {e}")

    return findings


def run_redaction_scan(dataset_name):
    """Thread target: scan all PDFs/DOCXs in a dataset for redactions."""
    global REDACT_PROGRESS, REDACT_RESULTS, REDACT_CANCEL

    folder    = Path(DATASETS_DIR) / dataset_name
    scan_exts = {'.pdf', '.docx', '.doc'}
    files     = sorted([f for f in folder.rglob("*")
                        if f.is_file() and not f.name.startswith("_")
                        and f.suffix.lower() in scan_exts])

    with REDACT_LOCK:
        REDACT_PROGRESS = {"status": "scanning", "done": 0,
                           "total": len(files), "current": "",
                           "dataset": dataset_name}
        REDACT_RESULTS  = []
        REDACT_CANCEL   = False

    results = []

    for i, fpath in enumerate(files):
        with REDACT_LOCK:
            if REDACT_CANCEL:
                REDACT_PROGRESS["status"] = "cancelled"
                REDACT_RESULTS = results
                return
            REDACT_PROGRESS["done"]    = i
            REDACT_PROGRESS["current"] = fpath.name

        try:
            ext = fpath.suffix.lower()
            findings = (detect_redactions_in_pdf(str(fpath))
                        if ext == '.pdf'
                        else detect_redactions_in_docx(str(fpath)))

            # Deduplicate: when a black bar and white box share the same location
            # (black bar drawn over white box), keep only the black bar
            if findings:
                black_locs = {f["location"] for f in findings if f["type"] == "Black Redaction Bar"}
                findings = [f for f in findings
                            if not (f["type"] == "White Box (Whiteout)" and f["location"] in black_locs)]

            if findings:
                pages      = sorted({f['page'] for f in findings if f.get('page')})
                type_counts = {}
                for f in findings:
                    type_counts[f['type']] = type_counts.get(f['type'], 0) + 1

                results.append({
                    'filename':        str(fpath.relative_to(folder)),
                    'title':           re.sub(r'[_\-]+', ' ', fpath.stem).strip(),
                    'ext':             ext,
                    'redaction_count': len(findings),
                    'pages':           pages,
                    'type_summary':    type_counts,
                    'findings':        findings,
                })

            with REDACT_LOCK:
                REDACT_RESULTS = results.copy()

        except Exception as e:
            print(f"  Redaction scan error {fpath.name}: {e}")

    with REDACT_LOCK:
        REDACT_PROGRESS["status"] = "done"
        REDACT_PROGRESS["done"]   = len(files)
        REDACT_RESULTS = results


def generate_redaction_report(results, dataset_name):
    """Generate Excel + PDF redaction report. Returns (xlsx_path, pdf_path)."""
    from datetime import datetime as dt
    timestamp = dt.now().strftime("%Y%m%d_%H%M%S")
    now_str   = dt.now().strftime("%B %d, %Y %I:%M %p")
    xlsx_path = os.path.join(EXPORTS_DIR, f"redaction_report_{dataset_name}_{timestamp}.xlsx")
    pdf_path  = os.path.join(EXPORTS_DIR, f"redaction_report_{dataset_name}_{timestamp}.pdf")

    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment

    hdr_fill  = PatternFill("solid", fgColor="1A1A2E")
    hdr_font  = Font(bold=True, color="FFFFFF", size=10)
    flag_fill = PatternFill("solid", fgColor="3D1010")
    stripe    = PatternFill("solid", fgColor="FFF0F0")
    wrap      = Alignment(wrap_text=True, vertical="top")
    center    = Alignment(horizontal="center", vertical="top")

    wb  = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "Summary"

    for col, (h, w) in enumerate([("Document", 52), ("Redaction Count", 18),
                                   ("Pages Affected", 30), ("Types Found", 55)], 1):
        cell = ws1.cell(row=1, column=col, value=h)
        cell.font = hdr_font; cell.fill = hdr_fill; cell.alignment = center
        ws1.column_dimensions[cell.column_letter].width = w
    ws1.freeze_panes = "A2"

    for row_idx, r in enumerate(results, 2):
        pages_str = ", ".join(f"p.{p}" for p in r['pages'][:20])
        if len(r['pages']) > 20:
            pages_str += f" (+{len(r['pages'])-20} more)"
        types_str = "; ".join(f"{t} ({c})" for t, c in r['type_summary'].items())
        ws1.cell(row=row_idx, column=1, value=r['filename']).alignment = wrap
        ws1.cell(row=row_idx, column=2, value=r['redaction_count']).alignment = center
        ws1.cell(row=row_idx, column=3, value=pages_str).alignment = wrap
        ws1.cell(row=row_idx, column=4, value=types_str).alignment = wrap
        fill = flag_fill if row_idx % 2 == 0 else stripe
        for col in range(1, 5):
            ws1.cell(row=row_idx, column=col).fill = fill
        ws1.row_dimensions[row_idx].height = 28

    ws2 = wb.create_sheet("Detailed Findings")
    for col, (h, w) in enumerate([("Document", 45), ("Page", 8),
                                   ("Type", 28), ("Location", 38), ("Details", 55)], 1):
        cell = ws2.cell(row=1, column=col, value=h)
        cell.font = hdr_font; cell.fill = hdr_fill; cell.alignment = center
        ws2.column_dimensions[cell.column_letter].width = w
    ws2.freeze_panes = "A2"

    row_idx = 2
    for r in results:
        for f in r['findings']:
            ws2.cell(row=row_idx, column=1, value=r['filename']).alignment = wrap
            ws2.cell(row=row_idx, column=2, value=f.get('page', '—')).alignment = center
            ws2.cell(row=row_idx, column=3, value=f['type']).alignment = wrap
            ws2.cell(row=row_idx, column=4, value=f['location']).alignment = wrap
            ws2.cell(row=row_idx, column=5, value=f['details']).alignment = wrap
            if row_idx % 2 == 0:
                for col in range(1, 6):
                    ws2.cell(row=row_idx, column=col).fill = stripe
            ws2.row_dimensions[row_idx].height = 24
            row_idx += 1

    wb.save(xlsx_path)

    # PDF report
    pdf_path_out = None
    try:
        from reportlab.lib.pagesizes import letter, landscape
        from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                        HRFlowable, Table, TableStyle)
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib import colors as rc
        from reportlab.lib.units import inch

        _pfx  = f"RD{timestamp}_"
        base  = getSampleStyleSheet()
        rdoc  = SimpleDocTemplate(pdf_path, pagesize=landscape(letter),
                                  leftMargin=0.75*inch, rightMargin=0.75*inch,
                                  topMargin=0.75*inch,  bottomMargin=0.75*inch)

        DARK   = rc.HexColor("#1a1a2e")
        ACCENT = rc.HexColor("#f87171")
        STRIPE = rc.HexColor("#fff5f5")

        title_s = ParagraphStyle(_pfx+"title", parent=base["Title"],   fontSize=18, spaceAfter=4)
        sub_s   = ParagraphStyle(_pfx+"sub",   parent=base["Normal"],  fontSize=9,
                                  textColor=rc.grey, spaceAfter=16)
        head_s  = ParagraphStyle(_pfx+"head",  parent=base["Heading1"],fontSize=13, spaceAfter=8)
        cell_s  = ParagraphStyle(_pfx+"cell",  parent=base["Normal"],  fontSize=8,  leading=11)

        def esc(t):
            return str(t).replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")

        total_redactions = sum(r['redaction_count'] for r in results)
        story = [
            Paragraph(esc(f"Redaction Report — {dataset_name}"), title_s),
            Paragraph(esc(f"Generated: {now_str}  ·  {len(results)} doc(s) with redactions  ·  "
                          f"{total_redactions} total instance(s)"), sub_s),
            HRFlowable(width="100%", thickness=2, color=ACCENT, spaceAfter=16),
            Paragraph("Summary", head_s),
        ]

        sum_data = [["Document", "Count", "Pages", "Types"]]
        for r in results:
            pages_short = ", ".join(f"p.{p}" for p in r['pages'][:8])
            if len(r['pages']) > 8: pages_short += "…"
            types_short = "; ".join(f"{t}({c})" for t, c in list(r['type_summary'].items())[:2])
            sum_data.append([
                Paragraph(esc(r['filename']), cell_s),
                str(r['redaction_count']),
                Paragraph(esc(pages_short), cell_s),
                Paragraph(esc(types_short), cell_s),
            ])
        if len(sum_data) == 1:
            sum_data.append(["No redactions found", "", "", ""])

        sum_tbl = Table(sum_data, colWidths=[3.5*inch, 0.7*inch, 2.2*inch, 3.1*inch])
        sum_tbl.setStyle(TableStyle([
            ("BACKGROUND",     (0,0), (-1,0), DARK),
            ("TEXTCOLOR",      (0,0), (-1,0), rc.white),
            ("FONTNAME",       (0,0), (-1,0), "Helvetica-Bold"),
            ("FONTSIZE",       (0,0), (-1,0), 9),
            ("ALIGN",          (1,0), (1,-1), "CENTER"),
            ("ROWBACKGROUNDS", (0,1), (-1,-1), [rc.white, STRIPE]),
            ("GRID",           (0,0), (-1,-1), 0.25, rc.HexColor("#ddcccc")),
            ("LINEBELOW",      (0,0), (-1,0), 1.5, ACCENT),
            ("VALIGN",         (0,0), (-1,-1), "MIDDLE"),
            ("TOPPADDING",     (0,0), (-1,-1), 4),
            ("BOTTOMPADDING",  (0,0), (-1,-1), 4),
        ]))
        story += [sum_tbl, Spacer(1, 20), Paragraph("Detailed Findings", head_s)]

        det_data = [["Document", "Pg", "Type", "Location", "Details"]]
        for r in results:
            for f in r['findings']:
                det_data.append([
                    Paragraph(esc(r['filename']), cell_s),
                    str(f.get('page', '—')),
                    Paragraph(esc(f['type']), cell_s),
                    Paragraph(esc(f['location']), cell_s),
                    Paragraph(esc(f['details']), cell_s),
                ])
        if len(det_data) == 1:
            det_data.append(["No findings", "", "", "", ""])

        det_tbl = Table(det_data, colWidths=[2.4*inch, 0.4*inch, 1.8*inch, 2.4*inch, 2.5*inch])
        det_tbl.setStyle(TableStyle([
            ("BACKGROUND",     (0,0), (-1,0), DARK),
            ("BACKGROUND",     (0,0), (-1,0), DARK),
            ("TEXTCOLOR",      (0,0), (-1,0), rc.white),
            ("FONTNAME",       (0,0), (-1,0), "Helvetica-Bold"),
            ("FONTSIZE",       (0,0), (-1,0), 9),
            ("ALIGN",          (1,0), (1,-1), "CENTER"),
            ("ROWBACKGROUNDS", (0,1), (-1,-1), [rc.white, STRIPE]),
            ("GRID",           (0,0), (-1,-1), 0.25, rc.HexColor("#ddcccc")),
            ("LINEBELOW",      (0,0), (-1,0), 1.5, ACCENT),
            ("VALIGN",         (0,0), (-1,-1), "MIDDLE"),
            ("TOPPADDING",     (0,0), (-1,-1), 4),
            ("BOTTOMPADDING",  (0,0), (-1,-1), 4),
        ]))
        story.append(det_tbl)
        rdoc.build(story)
        pdf_path_out = pdf_path
    except Exception as e:
        print(f"  Redaction PDF report error: {e}")

    return xlsx_path, pdf_path_out


# ── REDACTION ROUTES ──────────────────────────────────────────────────────────

@app.route("/redactions")
def redactions_page():
    return render_template("redactions.html")

@app.route("/redactions/scan", methods=["POST"])
def start_redaction_scan():
    data    = request.json or {}
    dataset = data.get("dataset", "").strip()
    if not dataset:
        return jsonify({"error": "dataset required"}), 400
    if not (Path(DATASETS_DIR) / dataset).exists():
        return jsonify({"error": "Dataset not found"}), 404
    t = threading.Thread(target=run_redaction_scan, args=(dataset,), daemon=True)
    t.start()
    return jsonify({"status": "started"})

@app.route("/redactions/status")
def redaction_status():
    with REDACT_LOCK:
        prog = {**REDACT_PROGRESS, "results_count": len(REDACT_RESULTS)}
    done  = prog.get("done",  0)
    total = prog.get("total", 1) or 1
    prog["percent"]      = int(done / total * 100)
    prog["message"]      = f"Scanning file {done} of {total}…"
    prog["current_file"] = prog.get("current", "")
    return jsonify(prog)

@app.route("/redactions/results")
def redaction_results_route():
    with REDACT_LOCK:
        results = list(REDACT_RESULTS)
    total_findings = sum(r.get("redaction_count", 0) for r in results)
    return jsonify({"documents": results, "total_findings": total_findings})

@app.route("/redactions/cancel", methods=["POST"])
def cancel_redaction():
    global REDACT_CANCEL
    with REDACT_LOCK:
        REDACT_CANCEL = True
    return jsonify({"ok": True})

@app.route("/redactions/export_zip", methods=["POST"])
def redaction_export_zip():
    import zipfile, io
    data    = request.json or {}
    dataset = data.get("dataset", "").strip()
    items   = data.get("items", [])
    if not dataset:
        with REDACT_LOCK:
            dataset = REDACT_PROGRESS.get("dataset", "")
    if not dataset:
        return jsonify({"error": "dataset required"}), 400
    if not items:
        with REDACT_LOCK:
            items = [r["filename"] for r in REDACT_RESULTS]
    folder  = Path(DATASETS_DIR) / dataset
    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for rel in items:
            src = folder / rel
            try:
                src.resolve().relative_to(folder.resolve())
            except ValueError:
                continue
            if src.exists():
                zf.write(str(src), arcname="redacted_docs/" + rel.replace("\\", "/"))
    zip_buf.seek(0)
    return send_file(zip_buf, mimetype="application/zip", as_attachment=True,
                     download_name=f"redacted_docs_{dataset}.zip")

@app.route("/redactions/export_list", methods=["POST"])
def redaction_export_list():
    """Export flat findings list as Excel."""
    with REDACT_LOCK:
        results = list(REDACT_RESULTS)
        dataset = REDACT_PROGRESS.get("dataset", "export")
    if not results:
        return jsonify({"error": "No results to export"}), 400
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment
        from datetime import datetime as dt
        timestamp = dt.now().strftime("%Y%m%d_%H%M%S")
        path      = os.path.join(EXPORTS_DIR, f"redaction_list_{dataset}_{timestamp}.xlsx")
        wb        = openpyxl.Workbook()
        ws        = wb.active
        ws.title  = "Findings"
        hdr_fill  = PatternFill("solid", fgColor="1A1A2E")
        hdr_font  = Font(bold=True, color="FFFFFF", size=10)
        stripe    = PatternFill("solid", fgColor="FFF0F0")
        _wrap     = Alignment(wrap_text=True, vertical="top")
        _ctr      = Alignment(horizontal="center", vertical="top")
        headers   = [("Document", 45), ("Page", 8), ("Type", 28),
                     ("Location", 38), ("Details", 55)]
        for col, (h, w) in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=h)
            cell.font = hdr_font; cell.fill = hdr_fill; cell.alignment = _ctr
            ws.column_dimensions[cell.column_letter].width = w
        ws.freeze_panes = "A2"
        row_idx = 2
        for r in results:
            for f in r.get("findings", []):
                ws.cell(row=row_idx, column=1, value=r["filename"]).alignment = _wrap
                ws.cell(row=row_idx, column=2, value=f.get("page", "—")).alignment = _ctr
                ws.cell(row=row_idx, column=3, value=f["type"]).alignment = _wrap
                ws.cell(row=row_idx, column=4, value=f["location"]).alignment = _wrap
                ws.cell(row=row_idx, column=5, value=f["details"]).alignment = _wrap
                if row_idx % 2 == 0:
                    for col in range(1, 6):
                        ws.cell(row=row_idx, column=col).fill = stripe
                ws.row_dimensions[row_idx].height = 24
                row_idx += 1
        wb.save(path)
        return send_file(path, as_attachment=True, download_name=os.path.basename(path))
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/redactions/export_report", methods=["POST", "GET"])
def redaction_export_report():
    data    = request.json or {}
    results = data.get("results", [])
    dataset = data.get("dataset", "")
    fmt     = request.args.get("format") or data.get("format", "xlsx")
    if not results:
        with REDACT_LOCK:
            results = list(REDACT_RESULTS)
    if not dataset:
        with REDACT_LOCK:
            dataset = REDACT_PROGRESS.get("dataset", "export")
    if not results:
        return jsonify({"error": "No results to export"}), 400
    try:
        xlsx_path, pdf_path = generate_redaction_report(results, dataset)
        if fmt == "pdf" and pdf_path:
            return send_file(pdf_path, as_attachment=True,
                             download_name=os.path.basename(pdf_path))
        return send_file(xlsx_path, as_attachment=True,
                         download_name=os.path.basename(xlsx_path))
    except Exception as e:
        return jsonify({"error": str(e)}), 500


if __name__ == "__main__":
    print("Starting Universal Search (Dev) at http://localhost:5000")
    print(f"Data directory: {BASE_DIR}")
    app.run(debug=False, port=5000)
